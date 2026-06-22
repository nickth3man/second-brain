"""Office document parser (§6).

Uses ``mammoth`` to convert document files, ``openpyxl`` for modern Excel
workbooks, and ``python-pptx`` for modern PowerPoint decks.
"""

from __future__ import annotations

from pathlib import Path

MAX_TABLE_COLUMNS = 40


async def parse_office(path: Path) -> str:
    """Convert an office document to markdown.

    Args:
        path: Path to the office file.

    Returns:
        Markdown text of the document content, or a clear unsupported/error note.
    """
    suffix = path.suffix.lower()

    if suffix in {".docx", ".doc", ".odt", ".rtf"}:
        try:
            import mammoth

            with open(path, "rb") as f:
                result = mammoth.convert_to_markdown(f)
            return result.value
        except Exception as e:
            return f"Office parse failed: {e}"

    if suffix == ".xlsx":
        return _parse_xlsx(path)

    if suffix == ".xls":
        return (
            "[legacy .xls parsing unsupported: binary Excel workbooks require "
            "an optional legacy parser/converter; save as .xlsx for safe "
            "value-only parsing]"
        )

    if suffix == ".pptx":
        return _parse_pptx(path)

    if suffix == ".ppt":
        return (
            "[legacy .ppt parsing unsupported: binary PowerPoint decks require "
            "an optional legacy converter; save as .pptx for markdown parsing]"
        )

    # Fallback for any unknown sub-extension within the office stage
    return f"[{suffix} office format unsupported]"


def _parse_xlsx(path: Path) -> str:
    """Convert a modern Excel workbook to markdown tables.

    ``data_only=True`` prevents formula execution and exposes only cached cell
    values already stored in the workbook by the spreadsheet application.
    """
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        return f"Spreadsheet parse failed: {exc}"

    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        visible_rows = [
            [_cell_to_markdown(value) for value in row[:MAX_TABLE_COLUMNS]]
            for row in rows
            if any(value is not None and str(value) != "" for value in row)
        ]
        sections.append(f"## Sheet: {sheet.title}")
        if not visible_rows:
            sections.append("_Empty sheet._")
            continue

        width = max(len(row) for row in visible_rows)
        normalized = [row + [""] * (width - len(row)) for row in visible_rows]
        headers = normalized[0]
        if not any(headers):
            headers = [f"Column {idx + 1}" for idx in range(width)]
            body_rows = normalized
        else:
            body_rows = normalized[1:]
        sections.append(_markdown_table(headers, body_rows))

    workbook.close()
    return "\n\n".join(sections)


def _parse_pptx(path: Path) -> str:
    """Convert a modern PowerPoint deck to structured markdown."""
    try:
        from pptx import Presentation

        deck = Presentation(path)
    except Exception as exc:
        return f"Presentation parse failed: {exc}"

    sections: list[str] = []
    for idx, slide in enumerate(deck.slides, 1):
        title = _slide_title(slide) or f"Slide {idx}"
        sections.append(f"## Slide {idx}: {title}")
        texts = _slide_texts(slide)
        if texts:
            sections.extend(f"- {text}" for text in texts)
        else:
            sections.append("_No slide text._")

        notes = _slide_notes(slide)
        if notes:
            sections.append("\n### Speaker Notes")
            sections.append(notes)
    return "\n\n".join(sections)


def _cell_to_markdown(value: object) -> str:
    if value is None:
        return ""
    return str(value).replace("\n", "<br>").replace("|", "\\|")


def _markdown_table(headers: list[str], rows: list[list[str]]) -> str:
    head = "| " + " | ".join(headers) + " |"
    sep = "| " + " | ".join("---" for _ in headers) + " |"
    body = ["| " + " | ".join(row) + " |" for row in rows]
    return "\n".join([head, sep, *body])


def _slide_title(slide) -> str:
    title = getattr(slide.shapes, "title", None)
    if title is not None and getattr(title, "has_text_frame", False):
        return title.text_frame.text.strip()
    return ""


def _slide_texts(slide) -> list[str]:
    text_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and getattr(shape.text_frame, "text", "").strip()
    ]
    text_shapes.sort(key=lambda shape: (shape.top, shape.left))
    seen_title = _slide_title(slide)
    texts: list[str] = []
    for shape in text_shapes:
        text = shape.text_frame.text.strip()
        if text == seen_title:
            continue
        texts.append(text.replace("\n", "\n  "))
    return texts


def _slide_notes(slide) -> str:
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""
    return notes
